"""
Le Centre F - Backend API pour l'Assistant IA Formation
========================================================
Architecture RAG (Retrieval-Augmented Generation) :
1. Ingestion : extraction texte depuis PDF/PPTX/DOCX
2. DÃ©coupage : chunks de 500-800 tokens avec overlap
3. Vectorisation : embeddings via sentence-transformers
4. Stockage : PostgreSQL + pgvector
5. Retrieval : recherche sÃ©mantique top-k
6. GÃ©nÃ©ration : API Claude avec contexte + sources
"""

import os
import json
import hashlib
from datetime import datetime, timedelta
from typing import Optional
from contextlib import asynccontextmanager

from fastapi import FastAPI, HTTPException, Depends, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from pydantic import BaseModel, Field
from pydantic_settings import BaseSettings


# ============================================================
# CONFIGURATION
# ============================================================

class Settings(BaseSettings):
    # API Claude (Anthropic)
    anthropic_api_key: str = "sk-ant-VOTRE-CLE-ICI"
    claude_model: str = "claude-sonnet-4-20250514"

    # Base de donnÃ©es PostgreSQL + pgvector
    database_url: str = "postgresql://centref:motdepasse@localhost:5432/centref_ai"

    # Embeddings
    embedding_model: str = "paraphrase-multilingual-MiniLM-L12-v2"
    chunk_size: int = 600
    chunk_overlap: int = 100
    top_k_results: int = 5

    # Auth
    secret_key: str = "CHANGEZ-MOI-EN-PRODUCTION-clef-secrete-64-chars"
    access_token_expire_minutes: int = 1440  # 24h

    class Config:
        env_file = ".env"

settings = Settings()


# ============================================================
# MODELES PYDANTIC
# ============================================================

class QuestionRequest(BaseModel):
    question: str = Field(..., min_length=3, max_length=2000)
    module_id: str = Field(..., pattern=r"^\d{3}$")
    conversation_id: Optional[str] = None

class QuestionResponse(BaseModel):
    answer: str
    sources: list
    conversation_id: str
    module_id: str
    processing_time_ms: int

class ModuleInfo(BaseModel):
    id: str
    name: str
    description: str
    document_count: int
    chunk_count: int

class IngestRequest(BaseModel):
    module_id: str
    file_path: str


# ============================================================
# MODULES DE FORMATION
# ============================================================

MODULES = {
    "001": {"name": "Plomb (CREP)", "description": "Constat de Risque d'Exposition au Plomb"},
    "002": {"name": "Amiante sans mention", "description": "Diagnostic amiante - niveau de base"},
    "003": {"name": "Amiante avec mention", "description": "Diagnostic amiante - niveau avancÃ©"},
    "004": {"name": "Ãnergie sans mention", "description": "DPE - Diagnostic de Performance ÃnergÃ©tique"},
    "005": {"name": "Ãnergie avec mention", "description": "DPE - niveau avancÃ© (tertiaire/ERP)"},
    "006": {"name": "Termites MÃ©tropole", "description": "Diagnostic termites France mÃ©tropolitaine"},
    "007": {"name": "ÃlectricitÃ©", "description": "Diagnostic installation Ã©lectrique"},
    "008": {"name": "Gaz", "description": "Diagnostic installation gaz"},
    "009": {"name": "Termites DROM", "description": "Diagnostic termites DOM-ROM"},
    "010": {"name": "DPEG", "description": "Diagnostic de Performance ÃnergÃ©tique Global"},
    "011": {"name": "DTG / PPT", "description": "Diagnostic Technique Global & Plan Pluriannuel de Travaux"},
}


# ============================================================
# EXTRACTION DE TEXTE
# ============================================================

def extract_text_from_pdf(file_path: str) -> list[dict]:
    """Extrait le texte page par page d'un PDF."""
    from PyPDF2 import PdfReader
    reader = PdfReader(file_path)
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            pages.append({
                "text": text.strip(),
                "page": i + 1,
                "source": os.path.basename(file_path),
                "type": "PDF"
            })
    return pages

def extract_text_from_pptx(file_path: str) -> list[dict]:
    """Extrait le texte slide par slide d'un PowerPoint."""
    from pptx import Presentation
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append({
                "text": "\n".join(texts),
                "page": i + 1,
                "source": os.path.basename(file_path),
                "type": "PPTX"
            })
    return slides

def extract_text_from_docx(file_path: str) -> list[dict]:
    """Extrait le texte d'un document Word."""
    from docx import Document
    doc = Document(file_path)
    paragraphs = []
    current_section = ""
    current_texts = []

    for para in doc.paragraphs:
        if para.style.name.startswith("Heading"):
            if current_texts:
                paragraphs.append({
                    "text": "\n".join(current_texts),
                    "page": len(paragraphs) + 1,
                    "source": os.path.basename(file_path),
                    "type": "DOCX",
                    "section": current_section
                })
            current_section = para.text.strip()
            current_texts = [current_section]
        elif para.text.strip():
            current_texts.append(para.text.strip())

    if current_texts:
        paragraphs.append({
            "text": "\n".join(current_texts),
            "page": len(paragraphs) + 1,
            "source": os.path.basename(file_path),
            "type": "DOCX",
            "section": current_section
        })
    return paragraphs


# ============================================================
# CHUNKING (DÃCOUPAGE)
# ============================================================

def chunk_documents(pages: list[dict], chunk_size: int = 600, overlap: int = 100) -> list[dict]:
    """DÃ©coupe les pages en chunks avec overlap pour le RAG."""
    from langchain_text_splitters import RecursiveCharacterTextSplitter

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
        length_function=len,
    )

    chunks = []
    for page in pages:
        texts = splitter.split_text(page["text"])
        for i, text in enumerate(texts):
            chunk_id = hashlib.md5(f"{page['source']}_{page['page']}_{i}".encode()).hexdigest()
            chunks.append({
                "id": chunk_id,
                "text": text,
                "source": page["source"],
                "page": page["page"],
                "type": page.get("type", ""),
                "section": page.get("section", ""),
                "chunk_index": i
            })
    return chunks


# ============================================================
# VECTORISATION ET STOCKAGE
# ============================================================

class VectorStore:
    """Gestion du stockage vectoriel avec PostgreSQL + pgvector."""

    def __init__(self, database_url: str, embedding_model: str):
        self.database_url = database_url
        self.model_name = embedding_model
        self._model = None

    @property
    def model(self):
        if self._model is None:
            from sentence_transformers import SentenceTransformer
            self._model = SentenceTransformer(self.model_name)
        return self._model

    def get_connection(self):
        import psycopg2
        return psycopg2.connect(self.database_url)

    def init_db(self):
        """CrÃ©e les tables nÃ©cessaires."""
        conn = self.get_connection()
        cur = conn.cursor()
        cur.execute("CREATE EXTENSION IF NOT EXISTS vector;")
        cur.execute("""
            CREATE TABLE IF NOT EXISTS documents (
                id TEXT PRIMARY KEY,
                module_id TEXT NOT NULL,
                filename TEXT NOT NULL,
                file_type TEXT,
                ingested_at TIMESTAMP DEFAULT NOW(),
                chunk_count INTEGER DEFAULT 0
            );
        """)
        cur.execute("""
            CREATE TABLE IF NOT EXISTS chunks (
                id TEXT PRIMARY KEY,
                document_id TEXT REFERENCES documents(id),
                module_id TEXT NOT NULL,
                content TEXT NOT NULL,
                source_file TEXT,
                page_number INTEGER,
                section TEXT,
                embedding vector(384),
                created_at TIMESTAMP DEFAULT NOW()
            );
        """)
        cur.execute("""
            CREATE INDEX IF NOT EXISTS idx_chunks_module ON chunks(module_id);
        """)
        cur.execute("""
            CREATE TABLE IF NOT EXISTS conversations (
                id TEXT PRIMARY KEY,
                user_id TEXT,
                module_id TEXT,
                created_at TIMESTAMP DEFAULT NOW()
            );
        """)
        cur.execute("""
            CREATE TABLE IF NOT EXISTS messages (
                id SERIAL PRIMARY KEY,
                conversation_id TEXT REFERENCES conversations(id),
                role TEXT NOT NULL,
                content TEXT NOT NULL,
                sources JSONB,
                created_at TIMESTAMP DEFAULT NOW()
            );
        """)
        conn.commit()
        cur.close()
        conn.close()

    def embed_text(self, text: str) -> list[float]:
        """GÃ©nÃ¨re l'embedding d'un texte."""
        return self.model.encode(text).tolist()

    def embed_batch(self, texts: list[str]) -> list[list[float]]:
        """GÃ©nÃ¨re les embeddings d'un batch de textes."""
        return self.model.encode(texts).tolist()

    def store_chunks(self, module_id: str, doc_id: str, filename: str, chunks: list[dict]):
        """Stocke les chunks avec leurs embeddings dans PostgreSQL."""
        conn = self.get_connection()
        cur = conn.cursor()

        # Enregistrer le document
        cur.execute(
            "INSERT INTO documents (id, module_id, filename, file_type, chunk_count) VALUES (%s, %s, %s, %s, %s) ON CONFLICT (id) DO UPDATE SET chunk_count = %s",
            (doc_id, module_id, filename, chunks[0].get("type", ""), len(chunks), len(chunks))
        )

        # GÃ©nÃ©rer les embeddings par batch
        texts = [c["text"] for c in chunks]
        embeddings = self.embed_batch(texts)

        # Stocker chaque chunk
        for chunk, emb in zip(chunks, embeddings):
            cur.execute(
                """INSERT INTO chunks (id, document_id, module_id, content, source_file, page_number, section, embedding)
                   VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
                   ON CONFLICT (id) DO NOTHING""",
                (chunk["id"], doc_id, module_id, chunk["text"], chunk["source"],
                 chunk["page"], chunk.get("section", ""), emb)
            )

        conn.commit()
        cur.close()
        conn.close()

    def search(self, query: str, module_id: str, top_k: int = 5) -> list[dict]:
        """Recherche sÃ©mantique des chunks les plus pertinents."""
        # VÃ©rifier d'abord s'il y a des chunks pour ce module
        # (Ã©vite de charger le modÃ¨le d'embeddings si aucun document n'est indexÃ©)
        try:
            conn = self.get_connection()
            cur = conn.cursor()
            cur.execute("SELECT COUNT(*) FROM chunks WHERE module_id = %s", (module_id,))
            count = cur.fetchone()[0]
            if count == 0:
                cur.close()
                conn.close()
                return []
        except Exception:
            # Connexion DB ou table non disponible - mode sans base
            return []

        query_embedding = self.embed_text(query)

        cur.execute(
            """SELECT content, source_file, page_number, section,
                      1 - (embedding <=> %s::vector) as similarity
               FROM chunks
               WHERE module_id = %s
               ORDER BY embedding <=> %s::vector
               LIMIT %s""",
            (query_embedding, module_id, query_embedding, top_k)
        )

        results = []
        for row in cur.fetchall():
            results.append({
                "content": row[0],
                "source": row[1],
                "page": row[2],
                "section": row[3],
                "similarity": round(float(row[4]), 4)
            })

        cur.close()
        conn.close()
        return results


# ============================================================
# GÃNÃRATION IA (Claude API)
# ============================================================

async def generate_answer(question: str, context_chunks: list[dict], module_name: str) -> dict:
    """GÃ©nÃ¨re une rÃ©ponse avec l'API Claude en mode RAG (ou connaissances gÃ©nÃ©rales si pas de documents)."""
    import anthropic

    client = anthropic.Anthropic(api_key=settings.anthropic_api_key)

    if context_chunks:
        # Mode RAG : rÃ©ponse basÃ©e sur les documents indexÃ©s
        context_parts = []
        for i, chunk in enumerate(context_chunks):
            source_info = f"[Source: {chunk['source']}, Page {chunk['page']}"
            if chunk.get('section'):
                source_info += f", Section: {chunk['section']}"
            source_info += f", Pertinence: {chunk['similarity']}]"
            context_parts.append(f"--- Extrait {i+1} {source_info} ---\n{chunk['content']}")

        context = "\n\n".join(context_parts)

        system_prompt = f"""Tu es l'assistant IA de formation du Centre F, spÃ©cialisÃ© dans les diagnostics immobiliers.
Tu rÃ©ponds aux questions des apprenants du module "{module_name}".

RÃGLES STRICTES :
1. RÃ©ponds UNIQUEMENT Ã  partir des extraits de documents fournis ci-dessous.
2. Si l'information n'est pas dans les extraits, dis-le clairement.
3. Cite TOUJOURS tes sources (nom du document, page, section) pour chaque affirmation.
4. Utilise un langage professionnel mais accessible.
5. Structure ta rÃ©ponse avec des paragraphes clairs.
6. Mets en gras les Ã©lÃ©ments clÃ©s avec **texte**.

EXTRAITS DES SUPPORTS DE FORMATION :
{context}"""
    else:
        # Mode connaissances gÃ©nÃ©rales : pas de documents indexÃ©s
        system_prompt = f"""Tu es l'assistant IA de formation du Centre F, spÃ©cialisÃ© dans les diagnostics immobiliers.
Tu rÃ©ponds aux questions des apprenants du module "{module_name}".

IMPORTANT : Les supports de formation de ce module n'ont pas encore Ã©tÃ© indexÃ©s dans la base de donnÃ©es.
Tu dois rÃ©pondre en utilisant tes connaissances gÃ©nÃ©rales sur le sujet du diagnostic immobilier.

RÃGLES :
1. RÃ©ponds de maniÃ¨re prÃ©cise et professionnelle en te basant sur la rÃ©glementation franÃ§aise en vigueur.
2. PrÃ©cise clairement que ta rÃ©ponse est basÃ©e sur tes connaissances gÃ©nÃ©rales et non sur les supports du Centre F.
3. Mentionne les textes rÃ©glementaires pertinents (arrÃªtÃ©s, normes NF, Code de la SantÃ© Publique, etc.).
4. Utilise un langage professionnel mais accessible.
5. Structure ta rÃ©ponse avec des paragraphes clairs.
6. Mets en gras les Ã©lÃ©ments clÃ©s avec **texte**."""

    message = client.messages.create(
        model=settings.claude_model,
        max_tokens=1500,
        system=system_prompt,
        messages=[{"role": "user", "content": question}]
    )

    answer_text = message.content[0].text

    # Extraire les sources utilisÃ©es (uniquement en mode RAG)
    sources = []
    if context_chunks:
        seen = set()
        for chunk in context_chunks:
            key = f"{chunk['source']}_{chunk['page']}"
            if key not in seen and chunk['similarity'] > 0.3:
                seen.add(key)
                source_type = "Support de formation"
                name = chunk['source'].lower()
                if "nf " in name or "norme" in name:
                    source_type = "Norme"
                elif "arrÃªtÃ©" in name or "dÃ©cret" in name or "arretÃ©" in name:
                    source_type = "RÃ©glementation"
                elif "code" in name or "loi" in name:
                    source_type = "Loi"

                sources.append({
                    "document": chunk['source'],
                    "page": chunk['page'],
                    "section": chunk.get('section', ''),
                    "type": source_type,
                    "relevance": chunk['similarity']
                })

    return {"answer": answer_text, "sources": sources[:5]}


# ============================================================
# APPLICATION FASTAPI
# ============================================================

@asynccontextmanager
async def lifespan(app: FastAPI):
    # Startup
    print("Le Centre F - Assistant IA Backend")
    print(f"ModÃ¨le IA : {settings.claude_model}")
    print(f"ModÃ¨le Embeddings : {settings.embedding_model}")
    yield
    # Shutdown
    print("ArrÃªt du serveur...")

app = FastAPI(
    title="Le Centre F - Assistant IA Formation",
    description="API backend pour l'assistant IA de formation aux diagnostics immobiliers",
    version="1.0.0",
    lifespan=lifespan
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # En production : restreindre Ã  votre domaine
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Instance du vector store
vector_store = VectorStore(settings.database_url, settings.embedding_model)


# ============================================================
# ENDPOINTS
# ============================================================

@app.get("/")
async def root():
    return {"service": "Le Centre F - Assistant IA", "status": "online", "version": "1.0.0"}

@app.get("/api/modules")
async def list_modules():
    """Liste tous les modules de formation disponibles."""
    return [{"id": k, **v} for k, v in MODULES.items()]

@app.post("/api/ask", response_model=QuestionResponse)
async def ask_question(req: QuestionRequest):
    """Pose une question Ã  l'assistant IA sur un module."""
    import time
    start = time.time()

    if req.module_id not in MODULES:
        raise HTTPException(status_code=400, detail=f"Module {req.module_id} inconnu")

    module = MODULES[req.module_id]

    # 1. Recherche sÃ©mantique des chunks pertinents
    chunks = vector_store.search(req.question, req.module_id, settings.top_k_results)

    # 2. GÃ©nÃ©ration de la rÃ©ponse avec Claude
    # Si pas de chunks indexÃ©s, Claude rÃ©pond avec ses connaissances gÃ©nÃ©rales
    result = await generate_answer(req.question, chunks, module["name"])

    elapsed = int((time.time() - start) * 1000)

    return QuestionResponse(
        answer=result["answer"],
        sources=result["sources"],
        conversation_id=req.conversation_id or hashlib.md5(str(time.time()).encode()).hexdigest()[:12],
        module_id=req.module_id,
        processing_time_ms=elapsed
    )

@app.post("/api/ingest")
async def ingest_document(req: IngestRequest):
    """IngÃ¨re un document dans la base vectorielle."""
    if not os.path.exists(req.file_path):
        raise HTTPException(status_code=404, detail="Fichier non trouvÃ©")

    ext = os.path.splitext(req.file_path)[1].lower()
    if ext == ".pdf":
        pages = extract_text_from_pdf(req.file_path)
    elif ext == ".pptx":
        pages = extract_text_from_pptx(req.file_path)
    elif ext == ".docx":
        pages = extract_text_from_docx(req.file_path)
    else:
        raise HTTPException(status_code=400, detail=f"Format {ext} non supportÃ©. Utilisez PDF, PPTX ou DOCX.")

    chunks = chunk_documents(pages, settings.chunk_size, settings.chunk_overlap)
    doc_id = hashlib.md5(req.file_path.encode()).hexdigest()

    vector_store.store_chunks(req.module_id, doc_id, os.path.basename(req.file_path), chunks)

    return {
        "status": "success",
        "document": os.path.basename(req.file_path),
        "module_id": req.module_id,
        "pages_extracted": len(pages),
        "chunks_created": len(chunks)
    }

@app.post("/api/ingest-module/{module_id}")
async def ingest_module(module_id: str, base_path: str):
    """IngÃ¨re automatiquement tous les documents d'un module."""
    if module_id not in MODULES:
        raise HTTPException(status_code=400, detail=f"Module {module_id} inconnu")

    results = []
    supported_ext = {".pdf", ".pptx", ".docx"}

    for root, dirs, files in os.walk(base_path):
        for f in files:
            if os.path.splitext(f)[1].lower() in supported_ext:
                file_path = os.path.join(root, f)
                try:
                    ext = os.path.splitext(f)[1].lower()
                    if ext == ".pdf":
                        pages = extract_text_from_pdf(file_path)
                    elif ext == ".pptx":
                        pages = extract_text_from_pptx(file_path)
                    elif ext == ".docx":
                        pages = extract_text_from_docx(file_path)
                    else:
                        continue

                    chunks = chunk_documents(pages)
                    doc_id = hashlib.md5(file_path.encode()).hexdigest()
                    vector_store.store_chunks(module_id, doc_id, f, chunks)

                    results.append({"file": f, "pages": len(pages), "chunks": len(chunks), "status": "ok"})
                except Exception as e:
                    results.append({"file": f, "status": "error", "error": str(e)})

    return {"module_id": module_id, "documents_processed": len(results), "details": results}

@app.get("/api/stats")
async def get_stats():
    """Statistiques globales de la base de connaissances."""
    try:
        conn = vector_store.get_connection()
        cur = conn.cursor()
        cur.execute("SELECT module_id, COUNT(*) FROM chunks GROUP BY module_id ORDER BY module_id")
        module_stats = {row[0]: row[1] for row in cur.fetchall()}
        cur.execute("SELECT COUNT(*) FROM documents")
        total_docs = cur.fetchone()[0]
        cur.execute("SELECT COUNT(*) FROM chunks")
        total_chunks = cur.fetchone()[0]
        cur.close()
        conn.close()
        return {"total_documents": total_docs, "total_chunks": total_chunks, "by_module": module_stats}
    except Exception:
        return {"total_documents": 0, "total_chunks": 0, "by_module": {}, "note": "Base de donnÃ©es non initialisÃ©e"}

@app.post("/api/init-db")
async def init_database():
    """Initialise la base de donnÃ©es (tables + extension pgvector)."""
    try:
        vector_store.init_db()
        return {"status": "success", "message": "Base de donnÃ©es initialisÃ©e avec succÃ¨s"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ============================================================
# POINT D'ENTRÃE
# ============================================================

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
